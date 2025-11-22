from pptx import Presentation
from llama_index.core import Document


def load_pptx(path, name):
    prs = Presentation(path)
    txt = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt.append(shape.text)

    return Document(
        text="\n".join(txt),
        metadata={"source": path}
    )
