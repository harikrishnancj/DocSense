
# backend/states/loader.py
"""
DocSense Loader V3 — Universal loader + visual analyzer + OCR + chart signals

Put this file at backend/states/loader.py and import Loader in your app_graph.
Returns LlamaIndex Document objects (from llama_index.core.Document) in state.documents.

Behaviors:
 - Loads PDFs (text + images + optional table extraction via camelot)
 - Extracts images from PDF/DOCX/PPTX and saves to disk
 - Runs vision analysis (GPT-4o style) per image -> caption + insights
 - OCR for scanned PDFs using PaddleOCR (preferred) or pytesseract fallback
 - Loads DOCX, TXT, CSV, Excel, PPTX, images, audio (placeholder transcription)
 - Extracts numeric/chart signals for use by the Visualizer state
 - Decorated with @traceable (LangSmith) so each load is traced
"""
'''
import os
import io
import re
import json
import tempfile
import zipfile
from typing import List, Dict, Tuple, Any

from langsmith import traceable
from llama_index.core import Document
from states.doc_state import DocState

# Common libs (ensure installed)
import pandas as pd
from PIL import Image
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

# OpenAI client — document uses 'OpenAI' as in your repo
from openai import OpenAI
client = OpenAI()  # requires OPENAI_API_KEY in env

# Optional libs: camelot for PDF tables, paddleocr for OCR
try:
    import camelot
    _HAS_CAMELOT = True
except Exception:
    _HAS_CAMELOT = False

# PaddleOCR preferred (accurate). Fallback to pytesseract if paddle not available.
_PADDLE_AVAILABLE = False
try:
    from paddleocr import PaddleOCR
    paddle_ocr = PaddleOCR(use_angle_cls=True, lang="en")
    _PADDLE_AVAILABLE = True
except Exception:
    _PADDLE_AVAILABLE = False

try:
    import pytesseract
    _PYTESSERACT_AVAILABLE = True
except Exception:
    _PYTESSERACT_AVAILABLE = False

# --- Helpers -----------------------------------------------------------------

EXTRACTED_IMG_DIR = "uploaded_docs/extracted_images"
os.makedirs(EXTRACTED_IMG_DIR, exist_ok=True)

def _save_pil_image(pil_img: Image.Image, filename_hint: str) -> str:
    """Save PIL image to extracted_images folder, return relative path."""
    os.makedirs(EXTRACTED_IMG_DIR, exist_ok=True)
    safe_name = re.sub(r"[^0-9A-Za-z._-]", "_", filename_hint)[:120]
    tmp_path = os.path.join(EXTRACTED_IMG_DIR, f"{safe_name}.png")
    pil_img.save(tmp_path, format="PNG")
    return tmp_path

def _run_ocr_on_pil(pil_img: Image.Image) -> str:
    """Run OCR using PaddleOCR if available else pytesseract. Return extracted text."""
    try:
        if _PADDLE_AVAILABLE:
            # PaddleOCR expects numpy array
            import numpy as np
            arr = np.array(pil_img.convert("RGB"))
            ocr_res = paddle_ocr.ocr(arr, cls=True)
            texts = [line[1][0] for page in ocr_res for line in page]
            return "\n".join(texts).strip()
        elif _PYTESSERACT_AVAILABLE:
            text = pytesseract.image_to_string(pil_img)
            return text.strip()
        else:
            return ""
    except Exception as e:
        return ""

def _extract_numeric_signals(text: str) -> Dict[str, Any]:
    """Extract years, percentages, currency numbers and plain numbers to help charting."""
    if not text:
        return {}
    years = re.findall(r"\b(19|20)\d{2}\b", text)
    years_full = re.findall(r"\b(19|20)\d{2}\b", text)  # same but kept for expansion
    perc = re.findall(r"\b\d+(?:\.\d+)?%\b", text)
    # currency: $1,234.56 or ₹1,234 or 1,234 USD etc.
    currency = re.findall(r"[\$\£\€\₹]\s?\d{1,3}(?:[,\d{3}]*)(?:\.\d+)?", text)
    numbers = re.findall(r"\b\d{1,3}(?:[,\d{3}]+)*(?:\.\d+)?\b", text)
    return {
        "years": years,
        "percentages": perc,
        "currency": currency,
        "numbers": numbers
    }

def analyze_image_with_lvm(pil_image: Image.Image) -> Tuple[str, str]:
    """
    Two-pass image analysis using GPT-4o-like API:
     - descriptive caption
     - insights (2-3 bullet observations)
    Returns (caption, insights)
    """
    try:
        buf = io.BytesIO()
        pil_image.save(buf, format="PNG")
        b64 = buf.getvalue()
        b64_data_url = f"data:image/png;base64,{b64.hex()}"  # fallback url-like token

        # NOTE: Many OpenAI wrappers support sending data URLs or base64 image data. If your
        # OpenAI client expects a different param, adapt the code here.
        # We'll call chat.completions.create with an "image_url" content entry similar to earlier examples.

        # Pass 1: caption
        try:
            resp1 = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "Describe this image in one concise sentence."},
                            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64.hex()}"}}
                        ],
                    }
                ],
                max_tokens=180,
            )
            caption = resp1.choices[0].message.content.strip()
        except Exception:
            # graceful fallback: ask model to describe from OCR if vision call fails
            caption = "[No caption generated]"

        # Pass 2: insight
        try:
            resp2 = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "Give 2-3 short insights or observations about this image."},
                            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64.hex()}"}}
                        ],
                    }
                ],
                max_tokens=250,
            )
            insights = resp2.choices[0].message.content.strip()
        except Exception:
            insights = ""

        return caption, insights
    except Exception as e:
        return f"[VisionError] {e}", ""

# --- file-type loaders ------------------------------------------------------

def load_pdf(path: str, filename: str) -> List[Document]:
    """Extract PDF pages text, images, run vision on images, optional tables via camelot."""
    docs: List[Document] = []
    try:
        pdf = fitz.open(path)
    except Exception as e:
        print(f"Failed to open PDF {filename}: {e}")
        return docs

    full_text = ""
    for page_num in range(len(pdf)):
        try:
            page = pdf[page_num]
            page_text = page.get_text()
            full_text += f"\n--- Page {page_num+1} ---\n{page_text}"
        except Exception as e:
            print(f"PDF page read error {filename} page {page_num+1}: {e}")
            continue

        # Extract images
        images_info = []
        try:
            images_info = page.get_images(full=True)
        except Exception:
            images_info = []

        for img_index, img in enumerate(images_info):
            try:
                xref = img[0]
                base = pdf.extract_image(xref)
                image_bytes = base.get("image")
                if not image_bytes:
                    continue
                pil_img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                # Save file and run vision
                img_path = _save_pil_image(pil_img, f"{filename}_p{page_num+1}_i{img_index}")
                caption, insights = analyze_image_with_lvm(pil_img)
                ocr_text = _run_ocr_on_pil(pil_img)

                docs.append(Document(
                    text=f"[PDF IMAGE]\nFilename: {filename}\nPage:{page_num+1}\nImageIndex:{img_index}\nCaption:{caption}\nInsights:{insights}\nOCR:{ocr_text}",
                    metadata={
                        "filename": filename,
                        "type": "pdf-image",
                        "page": page_num+1,
                        "image_index": img_index,
                        "image_path": img_path,
                        "caption": caption,
                        "insights": insights,
                        "ocr": ocr_text
                    }
                ))
            except Exception as e:
                print(f"Warning: image extraction failed for {filename} page {page_num+1} img {img_index}: {e}")
                continue

    # Optional table extraction via camelot
    if _HAS_CAMELOT:
        try:
            tables_texts = []
            tlist = camelot.read_pdf(path, pages="all", flavor="lattice")
            for t in tlist:
                if not t.df.empty:
                    tables_texts.append(t.df.to_string())
            if not tables_texts:
                tlist = camelot.read_pdf(path, pages="all", flavor="stream")
                for t in tlist:
                    if not t.df.empty:
                        tables_texts.append(t.df.to_string())
            for idx, tbl in enumerate(tables_texts):
                docs.append(Document(
                    text=f"[PDF TABLE]\nFilename:{filename}\nTableIndex:{idx}\n{tbl}",
                    metadata={"filename": filename, "type": "pdf-table", "table_index": idx}
                ))
        except Exception as e:
            print(f"Camelot extraction failed for {filename}: {e}")

    # Add full text doc
    docs.append(Document(text=full_text or "[NO_TEXT]", metadata={"filename": filename, "type": "pdf-text"}))
    pdf.close()
    return docs

def load_docx(path: str, filename: str) -> List[Document]:
    try:
        doc = DocxDocument(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        full_text = "\n".join(paragraphs).strip() or "[EMPTY DOCX]"
        docs: List[Document] = [Document(text=full_text, metadata={"filename": filename, "type":"docx"})]

        # extract inline images (if any)
        # python-docx doesn't easily expose raw images; fallback: try to read docx media by unzip
        try:
            with zipfile.ZipFile(path) as zf:
                for name in zf.namelist():
                    if name.startswith("word/media/") and not name.endswith("/"):
                        img_bytes = zf.read(name)
                        pil_img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                        img_path = _save_pil_image(pil_img, f"{filename}_docx_{os.path.basename(name)}")
                        caption, insights = analyze_image_with_lvm(pil_img)
                        ocr_text = _run_ocr_on_pil(pil_img)
                        docs.append(Document(
                            text=f"[DOCX IMAGE]\nFilename:{filename}\nImageName:{name}\nCaption:{caption}\nInsights:{insights}\nOCR:{ocr_text}",
                            metadata={"filename": filename, "type":"docx-image", "image_path": img_path}
                        ))
        except Exception:
            pass

        return docs
    except Exception as e:
        print(f"Failed to load DOCX {filename}: {e}")
        return []

def load_pptx(path: str, filename: str) -> List[Document]:
    try:
        prs = Presentation(path)
        slides_text = []
        docs: List[Document] = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                except Exception:
                    continue
            st = "\n".join([t for t in slide_text if t and t.strip()])
            slides_text.append(st)
            # extract images per slide (shapes with image)
            for shape in slide.shapes:
                try:
                    if shape.shape_type == 13: # Picture
                        img = shape.image
                        pil_img = Image.open(io.BytesIO(img.blob)).convert("RGB")
                        img_path = _save_pil_image(pil_img, f"{filename}_slide{i+1}")
                        caption, insights = analyze_image_with_lvm(pil_img)
                        docs.append(Document(
                            text=f"[PPTX IMAGE]\nFilename:{filename}\nSlide:{i+1}\nCaption:{caption}\nInsights:{insights}",
                            metadata={"filename": filename, "type":"pptx-image", "slide": i+1, "image_path": img_path}
                        ))
                except Exception:
                    continue

        docs.insert(0, Document(text="\n---SLIDES---\n".join(slides_text), metadata={"filename": filename, "type":"pptx"}))
        return docs
    except Exception as e:
        print(f"Failed to load PPTX {filename}: {e}")
        return []

def load_txt(path: str, filename: str) -> List[Document]:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return [Document(text=text or "[EMPTY TXT]", metadata={"filename": filename, "type":"txt"})]
    except Exception as e:
        print(f"Failed to load TXT {filename}: {e}")
        return []

def load_csv(path: str, filename: str) -> List[Document]:
    try:
        df = pd.read_csv(path)
        txt = df.to_string()
        return [Document(text=txt, metadata={"filename": filename, "type":"csv"})]
    except Exception as e:
        print(f"Failed to load CSV {filename}: {e}")
        return []

def load_excel(path: str, filename: str) -> List[Document]:
    try:
        sheets = pd.read_excel(path, sheet_name=None)
        docs = []
        for sheet_name, df in sheets.items():
            docs.append(Document(text=df.to_string(), metadata={"filename": filename, "type":"excel", "sheet":sheet_name}))
        return docs
    except Exception as e:
        print(f"Failed to load Excel {filename}: {e}")
        return []

def load_image(path: str, filename: str) -> List[Document]:
    try:
        pil_img = Image.open(path).convert("RGB")
        img_path = _save_pil_image(pil_img, f"{filename}")
        caption, insights = analyze_image_with_lvm(pil_img)
        ocr_text = _run_ocr_on_pil(pil_img)
        return [Document(text=f"[IMAGE]\nFilename:{filename}\nCaption:{caption}\nInsights:{insights}\nOCR:{ocr_text}",
                         metadata={"filename": filename, "type":"image", "image_path": img_path})]
    except Exception as e:
        print(f"Failed to load image {filename}: {e}")
        return []

def load_audio(path: str, filename: str) -> List[Document]:
    """
    Placeholder transcription: uses OpenAI speech-to-text if available (replace with your chosen transcriber).
    """
    try:
        # If using OpenAI whisper endpoint or other SDK, call here.
        # This is a simple placeholder that returns an empty transcription if not configured.
        transcription = ""
        # Example (pseudocode; adapt to your actual client if supported):
        # with open(path, "rb") as fh:
        #     resp = client.audio.transcriptions.create(file=fh, model="whisper-1")
        #     transcription = resp.text
        return [Document(text=f"[AUDIO]\nFilename:{filename}\nTranscription:{transcription}",
                         metadata={"filename": filename, "type":"audio"})]
    except Exception as e:
        print(f"Audio load/transcription failed for {filename}: {e}")
        return []

# --- Main Loader ------------------------------------------------------------

@traceable(name="loader")
def Loader(state: DocState) -> DocState:
    """
    Master loader: iterate files in state.folder_path, call file-specific loaders,
    append LlamaIndex Documents to state.documents, and enrich with numeric signals.
    """
    folder = state.folder_path or "uploaded_docs"
    os.makedirs(folder, exist_ok=True)
    all_docs: List[Document] = []

    for file in sorted(os.listdir(folder)):
        full_path = os.path.join(folder, file)
        if os.path.isdir(full_path):
            # skip directories (or optionally recurse)
            continue
        lower = file.lower()
        try:
            if lower.endswith(".pdf"):
                docs = load_pdf(full_path, file)
            elif lower.endswith(".docx"):
                docs = load_docx(full_path, file)
            elif lower.endswith((".pptx", ".ppt")):
                docs = load_pptx(full_path, file)
            elif lower.endswith(".txt"):
                docs = load_txt(full_path, file)
            elif lower.endswith(".csv"):
                docs = load_csv(full_path, file)
            elif lower.endswith((".xls", ".xlsx")):
                docs = load_excel(full_path, file)
            elif lower.endswith((".png", ".jpg", ".jpeg")):
                docs = load_image(full_path, file)
            elif lower.endswith((".mp3", ".wav", ".m4a")):
                docs = load_audio(full_path, file)
            elif lower.endswith(".zip"):
                # unzip and process extracted files (lightweight)
                try:
                    with zipfile.ZipFile(full_path, "r") as zf:
                        temp_dir = tempfile.mkdtemp()
                        zf.extractall(temp_dir)
                        # process extracted files
                        for sub in sorted(os.listdir(temp_dir)):
                            sp = os.path.join(temp_dir, sub)
                            # only one-level extraction for now
                            # call Loader recursively could cause state mutation complexity, so handle minimal types:
                            sub_lower = sub.lower()
                            if sub_lower.endswith(".pdf"):
                                docs = load_pdf(sp, sub)
                                all_docs.extend(docs)
                            elif sub_lower.endswith((".png", ".jpg", ".jpeg")):
                                docs = load_image(sp, sub)
                                all_docs.extend(docs)
                            elif sub_lower.endswith(".txt"):
                                docs = load_txt(sp, sub)
                                all_docs.extend(docs)
                        # cleanup optionally
                except Exception as e:
                    print(f"Zip processing failed for {file}: {e}")
                docs = []
            else:
                # Fallback: try to read as text
                docs = load_txt(full_path, file)
        except Exception as e:
            print(f"Error processing {file}: {e}")
            docs = []

        # enrich docs with numeric signals
        for d in docs:
            try:
                body = getattr(d, "text", "") or ""
                signals = _extract_numeric_signals(body)
                meta = d.metadata if isinstance(d.metadata, dict) else {}
                meta.setdefault("numeric_signals", signals)
                d.metadata = meta
            except Exception:
                pass
        all_docs.extend(docs)

    # Save to state
    state.documents = all_docs
    return state


'''
from langsmith import traceable
from llama_index.core import SimpleDirectoryReader
from states.doc_state import DocState
import os

@traceable(name="loader")
def Loader(state: DocState):
    """Load all documents from the folder into the state."""
    os.makedirs(state.folder_path, exist_ok=True)
    reader = SimpleDirectoryReader(state.folder_path)
    state.documents = reader.load_data()
    return state

'''from langsmith import traceable
from llama_index.core import Document
from states.doc_state import DocState
from model.vision import describe_image_gpt4o

import os
import fitz  # PyMuPDF
import pandas as pd
import pytesseract
from PIL import Image
from docx import Document as DocxReader
from pptx import Presentation

SUPPORTED_EXT = {
    ".pdf", ".txt", ".csv", ".xlsx", ".xls",
    ".docx", ".pptx", ".png", ".jpg", ".jpeg"
}

@traceable(name="loader")
def Loader(state: DocState):
    """Universal loader: PDFs, Excel, Images, CSV, DOCX, PPTX with table & image extraction."""
    
    folder = state.folder_path
    os.makedirs(folder, exist_ok=True)

    documents = []

    for filename in os.listdir(folder):
        path = os.path.join(folder, filename)
        ext = os.path.splitext(filename)[1].lower()

        if ext not in SUPPORTED_EXT:
            print(f"⚠ Skipping unsupported file: {filename}")
            continue

        if ext == ".pdf":
            documents.extend(load_pdf(path))

        elif ext in [".xlsx", ".xls"]:
            documents.extend(load_excel(path))

        elif ext == ".csv":
            documents.append(load_csv(path, filename))

        elif ext == ".txt":
            documents.append(load_txt(path, filename))

        elif ext == ".docx":
            documents.append(load_docx(path, filename))

        elif ext == ".pptx":
            documents.append(load_pptx(path, filename))

        elif ext in [".png", ".jpg", ".jpeg"]:
            documents.append(load_image(path, filename))

    state.documents = documents
    return state
'''
